import xlwings as xw
from pandas import DataFrame
from pptx import Presentation


class PPT:
    def __init__(self, input_file='template.pptx'):
        self.prs = Presentation(input_file)
        self.layouts = self.prs.slide_layouts
        self.slides = self.prs.slides

    def __getitem__(self, item):
        page_idx, shape_idx = item
        return self.slides[page_idx].shapes[shape_idx]

    def __setitem__(self, key, value):
        page_idx, shape_idx = key

        # 填入文字
        if type(value) is str:
            self.slides[page_idx].shapes[shape_idx].text = value
        # DataFrame填入表格
        elif type(value) is DataFrame:
            tb = self.slides[page_idx].shapes[shape_idx].table
            self.df2table(tb, value)
        # 异常
        else:
            raise ValueError(f"descriptor 'value' requires a 'str' or a 'DataFrame', but received a '{type(value)}'")

    def analyze_layouts(self, output_file='layouts_analyze.pptx'):
        # 遍历每个版式与占位符
        for s, layout in enumerate(self.prs.slide_layouts):
            slide = self.prs.slides.add_slide(layout)

            # 是否有标题占位符
            try:
                title = slide.shapes.title
                title.text = f'{s}样式-标题'
            except AttributeError:
                print(f'>>> page {i} has no title')

            # 将其他占位符(placeholders)命名为x样式x号
            for each in slide.placeholders:
                each.text = f'{s}样式-{each.placeholder_format.idx}号'

        # 保存
        self.save(output_file)

    def analyze_slides(self, output_file='slides_analyze.pptx'):
        # 遍历每页与每个shape
        for p, slide in enumerate(self.slides):
            for i, shape in enumerate(slide.shapes):
                shape.text = f'{p}页-{i}号'

        # 保存
        self.save(output_file)

    def save(self, output_file='ouput.pptx'):
        self.prs.save(output_file)
        print(f'>>> [{output_file}] saved.')

    def df2table(self, tb, df):
        # 确定表格行、列数
        rows, cols = df.shape

        # 填写表头
        columns = list(df.columns)
        for col, value in enumerate(columns):
            cell = tb.cell(0, col)
            cell.text = value

        # 填写数据
        df_matrix = df.as_matrix()
        for row in range(rows):
            for col in range(cols):
                value = df_matrix[row, col]
                cell = tb.cell(row + 1, col)
                cell.text = str(value)


class Excel:
    def __init__(self, visible=False):
        self.app = xw.App(visible=visible)
        self.wb = self.app.books[0]
        self.sheets = self.wb.sheets

    def __getitem__(self, item):
        return self.sheet_isin(item)

    def __setitem__(self, key, value):
        self.df2sheet(key, value)

    def sheet_isin(self, sheet_name):
        """不存就新建"""
        sht_names = [x.name for x in self.sheets]
        if sheet_name not in sht_names:
            sheet = self.sheets.add(sheet_name, after=len(self.sheets))
            print(f'>>> [{sheet_name}] added')
        else:
            sheet = self.sheets[sheet_name]
        return sheet

    def open_file(self, input_file):
        self.wb = self.app.books.open(input_file)
        print(f'>>> [{input_file}] is opened.')

    def df2sheet(self, sheet_name, df):
        sheet = self.sheet_isin(sheet_name)
        sheet.range('A1').value = df
        print(f'>>> DataFrame saved in [{sheet_name}].')

    def save(self, output_file='output.xlsx'):
        self.sheets['Sheet1'].delete()
        self.wb.save(output_file)
        self.wb.close()
        print(f'>>> [{output_file}] saved.')
